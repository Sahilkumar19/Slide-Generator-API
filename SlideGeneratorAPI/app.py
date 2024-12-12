from flask import Flask, request, jsonify, send_file
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import requests
import json
import os
import uuid
from datetime import datetime
from functools import wraps
import time

# Initialize Flask app
app = Flask(__name__)

# In-memory storage
presentations = {}
UPLOAD_FOLDER = 'presentations'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Configuration
GEMINI_API_KEY = "AIzaSyBwoukWCB3Tf4KrX_SVHQFvHkUwH10X4Cs"
GEMINI_API_URL = "https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash-latest:generateContent"

# Rate limiting decorator
def rate_limit(limit=100, window=3600):  # 100 requests per hour
    def decorator(f):
        requests_timestamps = []
        
        @wraps(f)
        def wrapped(*args, **kwargs):
            now = time.time()
            
            # Remove timestamps outside the window
            while requests_timestamps and requests_timestamps[0] < now - window:
                requests_timestamps.pop(0)
            
            if len(requests_timestamps) >= limit:
                return jsonify({'error': 'Rate limit exceeded'}), 429
            
            requests_timestamps.append(now)
            return f(*args, **kwargs)
        return wrapped
    return decorator

class PresentationGenerator:
    def __init__(self):
        self.layouts = {
            'title': 0,
            'bullet_points': 1,
            'two_column': 3,
            'content_with_image': 5
        }
    
    def generate_content(self, topic, num_slides=10):
        prompt = f"""Generate a {num_slides}-slide presentation for the topic '{topic}'.
        For each slide include:
        1. A header
        2. Content (50-60 words)
        3. Source citation
        Return as JSON array with keys 'header', 'content', and 'citation'."""

        headers = {"Content-Type": "application/json"}
        query_json = {
            "contents": [{"parts": [{"text": prompt}]}]
        }

        response = requests.post(
            GEMINI_API_URL,
            headers=headers,
            json=query_json,
            params={"key": GEMINI_API_KEY}
        )
        
        if response.status_code != 200:
            raise Exception("Content generation failed")

        result = response.json()
        content = result["candidates"][0]["content"]
        json_text = content["parts"][0]["text"].replace("```json\n", "").replace("\n```", "")
        return json.loads(json_text)

    def create_presentation(self, topic, config):
        prs = Presentation()
        
        # Apply theme/styling
        self.apply_theme(prs, config.get('theme', {}))
        
        # Generate content
        slides_content = self.generate_content(topic, config.get('num_slides', 10))
        
        for slide_content in slides_content:
            layout_type = config.get('layout', 'bullet_points')
            layout_index = self.layouts.get(layout_type, 1)
            
            slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
            
            # Add title
            if slide.shapes.title:
                title = slide.shapes.title
                title.text = slide_content['header']
            
            # Add content based on layout
            if layout_type == 'bullet_points':
                self.add_bullet_points(slide, slide_content)
            elif layout_type == 'two_column':
                self.add_two_column_content(slide, slide_content)
            elif layout_type == 'content_with_image':
                self.add_content_with_image(slide, slide_content)
            
            # Add citation
            if 'citation' in slide_content:
                self.add_citation(slide, slide_content['citation'])
        
        return prs

    def apply_theme(self, presentation, theme):
        # Apply background color
        for slide in presentation.slides:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
        
        # Apply font styles (would be expanded based on theme configuration)
        pass

    def add_bullet_points(self, slide, content):
        shapes = slide.shapes
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = content['content']
        
    def add_citation(self, slide, citation):
        left = Inches(0.5)
        top = Inches(6.5)
        width = Inches(9)
        height = Inches(0.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = f"Source: {citation}"
        p = tf.paragraphs[0]
        p.font.size = Pt(8)
        p.font.italic = True

@app.route('/api/v1/presentations', methods=['POST'])
@rate_limit()
def create_presentation():
    try:
        data = request.get_json()
        topic = data.get('topic')
        config = data.get('config', {})
        
        if not topic:
            return jsonify({'error': 'Topic is required'}), 400
            
        # Validate configuration
        if config.get('num_slides', 10) > 20:
            return jsonify({'error': 'Maximum 20 slides allowed'}), 400
        
        # Generate unique ID
        presentation_id = str(uuid.uuid4())
        
        # Create presentation
        generator = PresentationGenerator()
        prs = generator.create_presentation(topic, config)
        
        # Save presentation
        output_path = os.path.join(UPLOAD_FOLDER, f"{presentation_id}.pptx")
        prs.save(output_path)
        
        # Store presentation metadata
        presentations[presentation_id] = {
            'id': presentation_id,
            'topic': topic,
            'config': config,
            'created_at': datetime.utcnow().isoformat(),
            'file_path': output_path
        }
        
        return jsonify({
            'id': presentation_id,
            'message': 'Presentation created successfully'
        }), 201
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/v1/presentations/<presentation_id>', methods=['GET'])
@rate_limit()
def get_presentation(presentation_id):
    presentation = presentations.get(presentation_id)
    if not presentation:
        return jsonify({'error': 'Presentation not found'}), 404
    return jsonify(presentation), 200

@app.route('/api/v1/presentations/<presentation_id>/download', methods=['GET'])
@rate_limit()
def download_presentation(presentation_id):
    presentation = presentations.get(presentation_id)
    if not presentation:
        return jsonify({'error': 'Presentation not found'}), 404
    
    return send_file(
        presentation['file_path'],
        as_attachment=True,
        download_name=f"presentation_{presentation_id}.pptx"
    )

@app.route('/api/v1/presentations/<presentation_id>/configure', methods=['POST'])
@rate_limit()
def configure_presentation(presentation_id):
    presentation = presentations.get(presentation_id)
    if not presentation:
        return jsonify({'error': 'Presentation not found'}), 404
    
    config = request.get_json()
    
    # Update configuration
    presentation['config'].update(config)
    
    # Regenerate presentation with new config
    generator = PresentationGenerator()
    prs = generator.create_presentation(presentation['topic'], presentation['config'])
    prs.save(presentation['file_path'])
    
    return jsonify({
        'message': 'Presentation configuration updated successfully'
    }), 200

if __name__ == '__main__':
    app.run(debug=True)